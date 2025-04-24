from flask import Flask, request, jsonify, send_from_directory, send_file
from pptx import Presentation
from dotenv import load_dotenv
from google import genai
import subprocess
import base64
import os
import json
import time
import requests
import re
from pdf_processor import extract_text_from_pdf, save_extracted_pdf_text, save_enhanced_pdf_extraction, detect_math_content
import fitz  # PyMuPDF for more advanced PDF processing

# Load environment variables
load_dotenv()

app = Flask(__name__)
UPLOAD_FOLDER = 'slides'
ORIGINAL_FILES_FOLDER = 'original_files'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(ORIGINAL_FILES_FOLDER, exist_ok=True)

# Configure Gemini AI
os.environ["API_KEY"] = os.getenv("API_KEY")  # Ensure API_KEY is set in your .env file
client = genai.Client(api_key=os.environ["API_KEY"])

model = "gemini-2.0-flash"

# Stability AI API
STABILITY_API_KEY = os.getenv('STABILITY_API_KEY')
if not STABILITY_API_KEY:
    print("Warning: Missing STABILITY_API_KEY in .env - image generation will not work")

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    slides = []

    for idx, slide in enumerate(prs.slides, start=1):
        slide_data = {
            "slide_number": idx,
            "title": "",
            "content": [],
            "notes": "",
            "text": ""
        }
        
        # First pass: extract the title if it exists
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'text') and shape.text.strip():
                    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                        if shape.placeholder_format.type == 1:  # Title placeholder
                            slide_data["title"] = shape.text.strip()
                            break  # Found the title, exit the loop
            except AttributeError:
                continue
        
        # Second pass: extract all content including the title
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'text') and shape.text.strip():
                    # Add all text content
                    slide_data["content"].append(shape.text.strip())
            except AttributeError:
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_data["content"].append(shape.text.strip())
        
        # Get slide notes if they exist
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                slide_data["notes"] = slide.notes_slide.notes_text_frame.text.strip()
        except AttributeError:
            pass

        # Combine all text with title at the beginning
        slide_data["text"] = (
            (slide_data["title"] + "\n" if slide_data["title"] else "") +
            "\n".join(slide_data["content"])
        ).strip()

        slides.append(slide_data)

    return slides

def convert_ppt_to_pptx(ppt_path):
    try:
        # Full path to soffice.exe (change this if needed)
        soffice_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
        
        # Run the command and print the output for debugging
        result = subprocess.run([
            soffice_path, '--headless', '--convert-to', 'pptx',
            ppt_path, '--outdir', os.path.dirname(ppt_path)
        ], check=True, capture_output=True, text=True)
        
        # Print the output for debugging
        print("Conversion output:", result.stdout)
        print("Conversion error:", result.stderr)
        
        pptx_path = ppt_path.rsplit('.', 1)[0] + '.pptx'
        
        if os.path.exists(pptx_path):
            return pptx_path
        return None
    except Exception as e:
        print("Conversion failed:", e)
        return None

def save_extracted_text(slide_data, filename):
    # Create a JSON object that contains all slides
    presentation_data = {
        "filename": filename,
        "total_slides": len(slide_data),
        "extraction_time": time.strftime("%Y-%m-%d %H:%M:%S"),
        "slides": slide_data
    }
    
    # Save to a single JSON file
    with open(f'slides/{filename}_slides.json', 'w', encoding='utf-8') as f:
        json.dump(presentation_data, f, indent=2)

# Helper function to process slide references
def process_slide_references(text):
    # Process multiple patterns for slide ranges with different formats
    
    # Pattern: "Slides X-Y" (like "Slides 8-18" or "Slides 10-13")
    text = re.sub(
        r'Slides (\d+)-(\d+)', 
        r'<a href="#slide-range" data-range="\1-\2">Slides \1-\2</a>', 
        text
    )
    
    # Pattern for Pages instead of Slides
    text = re.sub(
        r'Pages (\d+)-(\d+)', 
        r'<a href="#slide-range" data-range="\1-\2">Pages \1-\2</a>', 
        text
    )
    
    # Pattern: "Slide X to Y" (handle another possible format)
    text = re.sub(
        r'Slide (\d+) to (\d+)', 
        r'<a href="#slide-range" data-range="\1-\2">Slides \1 to \2</a>', 
        text
    )
    
    # Pattern: "Page X to Y"
    text = re.sub(
        r'Page (\d+) to (\d+)', 
        r'<a href="#slide-range" data-range="\1-\2">Pages \1 to \2</a>', 
        text
    )
    
    # Pattern: "Slides X and Y" (not a range, but individual slides)
    text = re.sub(
        r'Slides (\d+) and (\d+)(?!\d)', 
        r'<a href="#slide-\1">Slide \1</a> and <a href="#slide-\2">Slide \2</a>', 
        text
    )
    
    # Pattern: "Pages X and Y"
    text = re.sub(
        r'Pages (\d+) and (\d+)(?!\d)', 
        r'<a href="#slide-\1">Page \1</a> and <a href="#slide-\2">Page \2</a>', 
        text
    )
    
    # Handle capitalization variations like "slide" instead of "Slide"
    text = re.sub(
        r'slide (\d+)', 
        r'<a href="#slide-\1">Slide \1</a>', 
        text,
        flags=re.IGNORECASE
    )
    
    # Handle "page" instead of "Page"
    text = re.sub(
        r'page (\d+)', 
        r'<a href="#slide-\1">Page \1</a>', 
        text,
        flags=re.IGNORECASE
    )

    # Process individual slide references - must come last
    text = re.sub(
        r'Slide (\d+)', 
        r'<a href="#slide-\1">Slide \1</a>', 
        text
    )
    
    # Process individual page references
    text = re.sub(
        r'Page (\d+)', 
        r'<a href="#slide-\1">Page \1</a>', 
        text
    )
    
    return text
        
# Update the call_gemini function to include visual elements from PDFs
def call_gemini(question, context, slide_number=None, include_visual_elements=True):
    try:
        slides = context.get("slides", [])
        has_visual_references = False
        visual_elements_context = ""
        image_data_to_include = []

        if slide_number is not None:
            # If a specific slide is selected, only use that slide's content
            slide_data = next((slide for slide in slides if slide["slide_number"] == slide_number), None)
            if slide_data:
                context_text = slide_data["text"]
                
                # Check if this is from an enhanced PDF with visual elements
                if include_visual_elements and context.get("images"):
                    # Get images for this page/slide
                    page_images = [img for img in context.get("images", []) if img["page"] == slide_number]
                    
                    if page_images:
                        has_visual_references = True
                        visual_elements_context = "\nVisual Elements on this page:\n"
                        
                        # Create detailed descriptions of each image
                        for i, image in enumerate(page_images):
                            image_desc = f"Image {i+1}: "
                            if "width" in image and "height" in image:
                                image_desc += f"Dimensions: {image['width']}x{image['height']}px. "
                            image_desc += f"Located on page {slide_number}. "
                            
                            # Store image data for direct inclusion
                            image_data_to_include.append({
                                "image_number": i+1,
                                "data_uri": image.get("data_uri", ""),
                                "description": image_desc
                            })
                            
                            visual_elements_context += f"- {image_desc}\n"
                
                # Also check for formulas
                if include_visual_elements and context.get("formulas"):
                    page_formulas = [f for f in context.get("formulas", []) if f["page"] == slide_number]
                    if page_formulas:
                        has_visual_references = True
                        if not visual_elements_context:
                            visual_elements_context = "\nVisual Elements on this page:\n"
                        
                        visual_elements_context += "Mathematical formulas:\n"
                        for i, formula in enumerate(page_formulas):
                            visual_elements_context += f"- Formula {i+1}: {formula['text']}\n"
                
                scope_notice = f"Answer based on content from Slide/Page {slide_number}:"
            else:
                context_text = ""
                scope_notice = f"No content found for Slide/Page {slide_number}."
        else:
            # If no specific slide is selected, use all slides' content
            context_text = "\n\n".join([f"Slide/Page {slide['slide_number']}:\n{slide['text']}" for slide in slides])
            
            # Add summary of visual elements for the whole document
            if include_visual_elements and (context.get("formulas") or context.get("images")):
                has_visual_references = True
                visual_elements_context = "\nVisual Elements Summary:\n"
                
                if context.get("formulas"):
                    formula_pages = set(f["page"] for f in context.get("formulas", []))
                    visual_elements_context += f"- Mathematical formulas found on pages: {', '.join(map(str, sorted(formula_pages)))}\n"
                
                if context.get("images"):
                    image_pages = set(img["page"] for img in context.get("images", []))
                    total_images = len(context.get("images", []))
                    visual_elements_context += f"- Total of {total_images} images found, distributed on pages: {', '.join(map(str, sorted(image_pages)))}\n"
                    
                    # Add details for each image
                    for i, image in enumerate(context.get("images", [])[:5]):  # Limit to first 5 images
                        image_desc = f"Image {i+1}: "
                        if "width" in image and "height" in image:
                            image_desc += f"Dimensions: {image['width']}x{image['height']}px. "
                        image_desc += f"Located on page {image.get('page', 'unknown')}. "
                        
                        # Store image data for direct inclusion
                        image_data_to_include.append({
                            "image_number": i+1,
                            "data_uri": image.get("data_uri", ""),
                            "description": image_desc
                        })
                        
                        visual_elements_context += f"- {image_desc}\n"
                    
            scope_notice = "Answer based on content from all slides/pages:"

        # Combine the context with visual elements
        full_context = context_text
        if has_visual_references:
            full_context += visual_elements_context

        # Prepare the multimodal content
        prompt_parts = [
            f"""As an AI tutor, please answer this question based on the slide/PDF content:

{scope_notice}
{full_context}

Question: {question}

Important formatting guidelines:
1. Do not use asterisks (*) for emphasis. Use HTML tags like <strong> or <em> instead.
2. When referencing individual slides or pages, use the format "Slide X" or "Page X".
3. When referencing a range of slides or pages, use the format "Slides X-Y" or "Pages X-Y".
4. Format any lists as proper HTML lists with <ul> and <li> tags.
5. Present your answer in clear, well-formatted paragraphs with proper spacing.
6. When referencing visual elements like formulas or images, be specific about their location.
7. If the question is about a formula, image, or other visual element, explicitly reference it in your answer."""
        ]
        
        # Include image data directly in the request
        for img_data in image_data_to_include:
            if img_data["data_uri"]:
                # Create multimodal content with both text and image
                # Format: prompt text, then image data
                prompt_parts.append({
                    "inlineData": {
                        "mimeType": "image/jpeg" if img_data["data_uri"].startswith("data:image/jpeg") else "image/png",
                        "data": img_data["data_uri"].split(',')[1]  # Extract base64 data
                    }
                })
                prompt_parts.append(f"This is {img_data['description']} Please analyze this image for the answer.")
        
        # Check if we have any image data
        has_images = len(image_data_to_include) > 0
        
        # Use different model for image analysis if needed
        selected_model = "gemini-1.5-pro" if has_images else model
        
        print(f"Using model: {selected_model}, Request has images: {has_images}")
        
        # Call appropriate Gemini API based on content
        if has_images:
            # Use multimodal generation for image content
            response = client.models.generate_content(
                model=selected_model,
                contents=prompt_parts
            )
        else:
            # Use text-only generation
            response = client.models.generate_content(
                model=model,
                contents=prompt_parts[0]  # Just the text prompt
            )

        processed_text = response.text
        
        # Replace any remaining asterisks for emphasis
        processed_text = re.sub(r'\*\*(.*?)\*\*', r'<strong>\1</strong>', processed_text)
        processed_text = re.sub(r'\*(.*?)\*', r'<em>\1</em>', processed_text)
        
        # Process links to slides and pages
        processed_text = process_slide_references(processed_text)
        
        return processed_text

    except Exception as e:
        print(f"Gemini API error: {str(e)}")
        return f"Error: Failed to get response from Gemini. {str(e)}"

# Update the ask_question route in app.py to handle presentation-specific queries with image support for math content
@app.route('/ask', methods=['POST'])
def ask_question():
    data = request.json
    question = data.get("question")
    slide_num = data.get("slide_number")  # This can be None if no specific slide is selected
    filename = data.get("filename")
    include_visual = data.get("include_visual_elements", True)  # Default to including visual elements

    try:
        # First check for enhanced PDF data
        enhanced_file_path = f'slides/{filename}_enhanced.json'
        standard_file_path = f'slides/{filename}_slides.json'
        
        print(f"Processing question about '{filename}', slide: {slide_num}, include_visual: {include_visual}")
        print(f"Looking for enhanced data at: {enhanced_file_path}")
        
        if os.path.exists(enhanced_file_path):
            # Use enhanced PDF data with visual elements
            with open(enhanced_file_path, 'r', encoding='utf-8') as f:
                presentation_data = json.load(f)
            
            # Check if we have math content on specific slides
            math_pages = []
            if "slides" in presentation_data:
                math_pages = [slide["slide_number"] for slide in presentation_data["slides"] 
                             if slide.get("has_math_content", False)]
                
            print(f"Detected math content on pages: {math_pages}")
                
            # Determine if we need to include page images in the query
            include_page_images = False
            page_images_to_include = []
            
            if include_visual:
                if slide_num is not None:
                    # Check if the specific slide has math content
                    slide_data = next((slide for slide in presentation_data.get("slides", []) 
                                     if slide["slide_number"] == slide_num), None)
                    if slide_data and slide_data.get("has_math_content", False) and slide_data.get("page_image"):
                        include_page_images = True
                        page_images_to_include.append({
                            "page": slide_num,
                            "image": slide_data["page_image"],
                            "description": f"Full page {slide_num} containing mathematical content"
                        })
                else:
                    # If searching all slides, include all math-containing pages (up to a reasonable limit)
                    for slide in presentation_data.get("slides", []):
                        if slide.get("has_math_content", False) and slide.get("page_image"):
                            include_page_images = True
                            # Limit to first 5 pages with math to keep request size reasonable
                            if len(page_images_to_include) < 5:
                                page_images_to_include.append({
                                    "page": slide["slide_number"],
                                    "image": slide["page_image"],
                                    "description": f"Full page {slide['slide_number']} containing mathematical content"
                                })
            
            # Call Gemini with the enhanced PDF data, including page images if needed
            response = call_gemini_with_math_support(question, presentation_data, slide_num, 
                                                     include_visual, page_images_to_include)
        elif os.path.exists(standard_file_path):
            # Fall back to standard text-only data
            print("Using standard text-only data")
            with open(standard_file_path, 'r', encoding='utf-8') as f:
                presentation_data = json.load(f)
                
            # Call Gemini with the presentation data (no page images in standard mode)
            response = call_gemini(question, presentation_data, slide_num, False)
        else:
            return jsonify({"error": f"Presentation '{filename}' not found"}), 404

        return jsonify({
            "question": question,
            "answer": response,
            "source_presentation": filename,
            "has_visual_elements": os.path.exists(enhanced_file_path)
        })
    except Exception as e:
        print(f"Error processing question: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

# New function to handle Gemini calls with math content page images
def call_gemini_with_math_support(question, context, slide_number=None, include_visual_elements=True, page_images=None):
    try:
        slides = context.get("slides", [])
        has_visual_references = False
        visual_elements_context = ""
        image_data_to_include = []

        if slide_number is not None:
            # If a specific slide is selected, only use that slide's content
            slide_data = next((slide for slide in slides if slide["slide_number"] == slide_number), None)
            if slide_data:
                context_text = slide_data["text"]
                
                # Check if this is a page with math content
                if slide_data.get("has_math_content", False) and include_visual_elements:
                    has_visual_references = True
                    visual_elements_context = "\nThis page contains mathematical content that may not be accurately represented as text.\n"
                
                # Check if this is from an enhanced PDF with visual elements
                if include_visual_elements and context.get("images"):
                    # Get images for this page/slide
                    page_images = [img for img in context.get("images", []) if img["page"] == slide_number]
                    
                    if page_images:
                        has_visual_references = True
                        if not visual_elements_context:
                            visual_elements_context = "\nVisual Elements on this page:\n"
                        
                        visual_elements_context += f"- {len(page_images)} images on page {slide_number}\n"
                        
                        # Store image data for direct inclusion (limit to first 3 for performance)
                        for i, image in enumerate(page_images[:3]):
                            image_desc = f"Image {i+1}: "
                            if "width" in image and "height" in image:
                                image_desc += f"Dimensions: {image['width']}x{image['height']}px. "
                            image_desc += f"Located on page {slide_number}. "
                            
                            # Store image data for direct inclusion
                            image_data_to_include.append({
                                "image_number": i+1,
                                "data_uri": image.get("data_uri", ""),
                                "description": image_desc
                            })
                
                # Also check for formulas
                if include_visual_elements and context.get("formulas"):
                    page_formulas = [f for f in context.get("formulas", []) if f["page"] == slide_number]
                    if page_formulas:
                        has_visual_references = True
                        if not visual_elements_context:
                            visual_elements_context = "\nVisual Elements on this page:\n"
                        
                        visual_elements_context += f"- {len(page_formulas)} mathematical formulas detected on page {slide_number}\n"
                
                scope_notice = f"Answer based on content from Slide/Page {slide_number}:"
            else:
                context_text = ""
                scope_notice = f"No content found for Slide/Page {slide_number}."
        else:
            # If no specific slide is selected, use all slides' content
            context_text = "\n\n".join([f"Slide/Page {slide['slide_number']}:\n{slide['text']}" for slide in slides])
            
            # Add summary of visual elements for the whole document
            if include_visual_elements:
                # Check for pages with math content
                math_pages = [slide["slide_number"] for slide in slides if slide.get("has_math_content", False)]
                
                if math_pages:
                    has_visual_references = True
                    visual_elements_context = "\nMathematical Content:\n"
                    visual_elements_context += f"- Mathematical notation detected on pages: {', '.join(map(str, sorted(math_pages)))}\n"
                
                # Add info about other visual elements
                if context.get("formulas") or context.get("images"):
                    has_visual_references = True
                    if not visual_elements_context:
                        visual_elements_context = "\nVisual Elements Summary:\n"
                    
                    if context.get("formulas"):
                        formula_pages = set(f["page"] for f in context.get("formulas", []))
                        visual_elements_context += f"- Mathematical formulas found on pages: {', '.join(map(str, sorted(formula_pages)))}\n"
                    
                    if context.get("images"):
                        image_pages = set(img["page"] for img in context.get("images", []))
                        total_images = len(context.get("images", []))
                        visual_elements_context += f"- Total of {total_images} images found, distributed on pages: {', '.join(map(str, sorted(image_pages)))}\n"
            
            scope_notice = "Answer based on content from all slides/pages:"

        # Combine the context with visual elements
        full_context = context_text
        if has_visual_references:
            full_context += visual_elements_context

        # Prepare the multimodal content
        prompt_parts = [
            f"""As an AI tutor, please answer this question based on the slide/PDF content:

{scope_notice}
{full_context}

Question: {question}

Important notes:
1. Some pages contain mathematical notation and formulas that might not be accurately represented as text.
2. I will provide images of pages with mathematical content to help you understand the notation correctly.
3. Please analyze both the text and the page images to provide an accurate answer.

Important formatting guidelines:
1. Do not use asterisks (*) for emphasis. Use HTML tags like <strong> or <em> instead.
2. When referencing individual slides or pages, use the format "Slide X" or "Page X".
3. When referencing a range of slides or pages, use the format "Slides X-Y" or "Pages X-Y".
4. Format any lists as proper HTML lists with <ul> and <li> tags.
5. Present your answer in clear, well-formatted paragraphs with proper spacing.
6. When referencing mathematical formulas, describe them accurately based on the page images."""
        ]
        
        # If we have specific page images for math content, prioritize those
        if page_images:
            for img_data in page_images:
                if "image" in img_data and img_data["image"]:
                    # Determine image MIME type
                    mime_type = "image/png"  # Default
                    if img_data["image"].startswith("data:image/jpeg"):
                        mime_type = "image/jpeg"
                    elif img_data["image"].startswith("data:image/png"):
                        mime_type = "image/png"
                        
                    # Extract base64 data
                    base64_data = img_data["image"].split(',')[1] if ',' in img_data["image"] else img_data["image"]
                    
                    # Add the image to prompt parts
                    prompt_parts.append({
                        "inlineData": {
                            "mimeType": mime_type,
                            "data": base64_data
                        }
                    })
                    prompt_parts.append(f"This is page {img_data['page']} containing mathematical content. Please analyze the mathematical notation in this image.")
        
        # Include other images if needed and we haven't already added too many
        elif len(image_data_to_include) > 0:
            for img_data in image_data_to_include[:3]:  # Limit to 3 images
                if img_data["data_uri"]:
                    # Create multimodal content with both text and image
                    # Format: prompt text, then image data
                    prompt_parts.append({
                        "inlineData": {
                            "mimeType": "image/jpeg" if img_data["data_uri"].startswith("data:image/jpeg") else "image/png",
                            "data": img_data["data_uri"].split(',')[1]  # Extract base64 data
                        }
                    })
                    prompt_parts.append(f"This is {img_data['description']} Please analyze this image for the answer.")
        
        # Check if we have any image data
        has_images = len(page_images) > 0 or len(image_data_to_include) > 0
        
        # Use different model for image analysis if needed
        selected_model = "gemini-1.5-pro" if has_images else model
        
        print(f"Using model: {selected_model}, Request has images: {has_images}")
        
        # Call appropriate Gemini API based on content
        if has_images:
            # Use multimodal generation for image content
            response = client.models.generate_content(
                model=selected_model,
                contents=prompt_parts
            )
        else:
            # Use text-only generation
            response = client.models.generate_content(
                model=model,
                contents=prompt_parts[0]  # Just the text prompt
            )

        processed_text = response.text
        
        # Replace any remaining asterisks for emphasis
        processed_text = re.sub(r'\*\*(.*?)\*\*', r'<strong>\1</strong>', processed_text)
        processed_text = re.sub(r'\*(.*?)\*', r'<em>\1</em>', processed_text)
        
        # Process links to slides and pages
        processed_text = process_slide_references(processed_text)
        
        return processed_text

    except Exception as e:
        print(f"Gemini API error: {str(e)}")
        return f"Error: Failed to get response from Gemini. {str(e)}"

def analyze_math_content_in_pdf(file_path):
    """
    Utility function to analyze PDF pages for mathematical content.
    Helps diagnose false positives and negatives in math content detection.
    """
    try:
        # Open the PDF with PyMuPDF
        doc = fitz.open(file_path)
        results = []
        
        print(f"Analyzing {file_path} for math content...")
        print(f"Total pages: {len(doc)}")
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            page_text = page.get_text()
            
            # Try different detection methods
            main_detection = detect_math_content(page_text)
            
            # Count math symbols
            math_symbols = r'[=+\-*/^√∫∑∏πλθ]'
            symbol_count = len(re.findall(math_symbols, page_text))
            
            # Look for Greek letters separately
            greek_letters = r'[αβγδεζηθικλμνξοπρστυφχψω]'
            greek_count = len(re.findall(greek_letters, page_text))
            
            # Look for equation patterns
            eq_patterns = r'[a-zA-Z]\s*=\s*[a-zA-Z0-9]'
            equation_count = len(re.findall(eq_patterns, page_text))
            
            # Check for specific formula patterns like λ = 2d Sin θ
            specific_formulas = any(re.search(r'λ\s*=\s*2d\s*Sin', page_text, re.IGNORECASE))
            
            # Get the first few lines to check if it's a title page
            first_lines = '\n'.join(page_text.split('\n')[:5])
            is_title_page = bool(re.match(r'^Chapter|^\d+\.\d+\s+[A-Z]', first_lines))
            
            # Examine blocks for potential formulas
            blocks = page.get_text("dict")["blocks"]
            block_math_count = 0
            for block in blocks:
                if "lines" in block:
                    for line in block["lines"]:
                        line_text = "".join([span["text"] for span in line["spans"]])
                        if detect_math_content(line_text):
                            block_math_count += 1
            
            # Compile the results
            page_result = {
                "page_number": page_num + 1,
                "math_detected": main_detection,
                "symbol_count": symbol_count,
                "greek_letter_count": greek_count,
                "equation_count": equation_count,
                "specific_formulas_found": specific_formulas,
                "is_title_page": is_title_page,
                "block_math_count": block_math_count,
                "first_lines": first_lines,
                "likely_has_math": (symbol_count > 3) or (greek_count > 0) or (equation_count > 0) or specific_formulas or (block_math_count > 1)
            }
            
            results.append(page_result)
            
            # Print a summary
            print(f"\nPage {page_num + 1}:")
            print(f"  Math detected: {main_detection}")
            print(f"  Symbol count: {symbol_count}")
            print(f"  Greek letters: {greek_count}")
            print(f"  Equations: {equation_count}")
            print(f"  Block math count: {block_math_count}")
            print(f"  First lines: {first_lines[:100]}...")
            print(f"  CONCLUSION: {'MATH CONTENT' if page_result['likely_has_math'] else 'NO MATH CONTENT'}")
        
        return results
    
    except Exception as e:
        print(f"Error analyzing math content: {str(e)}")
        return []

@app.route('/analyze-math/<path:filename>')
def analyze_math(filename):
    """API endpoint to analyze math content detection in a PDF file."""
    try:
        # Look for the file in original files folder
        file_path = os.path.join(ORIGINAL_FILES_FOLDER, filename)
        if not os.path.exists(file_path):
            # Try with .pdf extension
            file_path = os.path.join(ORIGINAL_FILES_FOLDER, f"{filename}.pdf")
            if not os.path.exists(file_path):
                return jsonify({"error": f"File not found: {filename}"}), 404
        
        # Use the math content analyzer from pdf_processor
        from pdf_processor import analyze_math_content_in_pdf
        results = analyze_math_content_in_pdf(file_path)
        
        return jsonify({
            "filename": filename,
            "total_pages": len(results),
            "analysis": results
        })
    
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

# Add endpoint to reprocess PDFs with the improved detection
@app.route('/reprocess-pdf/<path:filename>')
def reprocess_pdf(filename):
    """Reprocess a PDF file with improved math content detection."""
    try:
        # Verify the file exists
        file_path = os.path.join(ORIGINAL_FILES_FOLDER, filename)
        if not os.path.exists(file_path):
            # Try with .pdf extension
            file_path = os.path.join(ORIGINAL_FILES_FOLDER, f"{filename}.pdf")
            if not os.path.exists(file_path):
                return jsonify({"error": f"File not found: {filename}"}), 404
        
        # Get the base name without extension
        basename = os.path.splitext(os.path.basename(file_path))[0]
        
        # Import updated functions from pdf_processor
        from pdf_processor import save_enhanced_pdf_extraction
        
        # Reprocess the PDF with enhanced detection
        result = save_enhanced_pdf_extraction(file_path, basename)
        
        if result:
            return jsonify({
                "success": True,
                "filename": basename,
                "message": f"PDF reprocessed successfully with improved math detection",
                "math_pages": result.get("math_content_pages", [])
            })
        else:
            return jsonify({
                "success": False,
                "filename": basename,
                "message": "Failed to reprocess PDF"
            })
    
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

# This will help pinpoint where the PDF loading is failing
@app.route('/original-file/<path:filename>')
def serve_original_file(filename):
    """
    Serve the original file from the original files folder.
    Using path:filename to handle filenames with spaces.
    """
    print(f"Requested file: {filename}")
    print(f"Looking in folder: {ORIGINAL_FILES_FOLDER}")
    
    # Check if file exists without extension (the extension might be added by the frontend)
    base_path = os.path.join(ORIGINAL_FILES_FOLDER, filename)
    
    if os.path.exists(base_path):
        print(f"File found at: {base_path}")
        return send_file(base_path)
    
    # Try with different extension combinations
    if '.' not in filename:
        pdf_path = f"{base_path}.pdf"
        if os.path.exists(pdf_path):
            print(f"File found at: {pdf_path}")
            return send_file(pdf_path)
    
    # List available files for debugging
    available_files = os.listdir(ORIGINAL_FILES_FOLDER)
    print(f"Available files in {ORIGINAL_FILES_FOLDER}: {available_files}")
    
    return jsonify({"error": "File not found", "requested": filename}), 404

# Add a route to get visual elements from a PDF
@app.route('/pdf-visual-elements/<filename>/<int:page>')
def get_pdf_visual_elements(filename, page):
    enhanced_file_path = f'slides/{filename}_enhanced.json'
    
    if not os.path.exists(enhanced_file_path):
        return jsonify({"error": "Enhanced PDF data not found"}), 404
        
    try:
        with open(enhanced_file_path, 'r', encoding='utf-8') as f:
            pdf_data = json.load(f)
            
        # Extract formulas and images for the specified page
        formulas = [f for f in pdf_data.get("formulas", []) if f["page"] == page]
        images = [img for img in pdf_data.get("images", []) if img["page"] == page]
        
        return jsonify({
            "filename": filename,
            "page": page,
            "formulas": formulas,
            "images": images
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# Modify the upload route in app.py to ensure files are saved properly and paths are correctly tracked
@app.route('/upload', methods=['POST'])
def upload_slide():
    try:
        # Accept multiple files
        uploaded_files = request.files.getlist('files[]')
        
        if not uploaded_files:
            return jsonify({"error": "No files uploaded"}), 400
        
        presentations = []
        
        # Create folders if they don't exist
        os.makedirs(UPLOAD_FOLDER, exist_ok=True)
        os.makedirs(ORIGINAL_FILES_FOLDER, exist_ok=True)
        
        print(f"Upload folder: {os.path.abspath(UPLOAD_FOLDER)}")
        print(f"Original files folder: {os.path.abspath(ORIGINAL_FILES_FOLDER)}")

        for file in uploaded_files:
            if file.filename == '':
                continue  # Skip empty file inputs

            # Verify file extension - now include PDF
            if not file.filename.lower().endswith(('.ppt', '.pptx', '.pdf')):
                continue  # Skip invalid files
                
            print(f"Processing uploaded file: {file.filename}")

            # Save original file in the originals folder
            original_file_path = os.path.join(ORIGINAL_FILES_FOLDER, file.filename)
            file.save(original_file_path)
            print(f"Saved original file to: {original_file_path}")
            
            # Make sure the file exists where we expect it
            if not os.path.exists(original_file_path):
                print(f"WARNING: File not found at expected path: {original_file_path}")
            else:
                print(f"Confirmed file exists at: {original_file_path}")
                print(f"File size: {os.path.getsize(original_file_path)} bytes")
            
            # Also save in the slides folder for backward compatibility
            slides_file_path = os.path.join(UPLOAD_FOLDER, file.filename)
            with open(original_file_path, 'rb') as f_src:
                with open(slides_file_path, 'wb') as f_dst:
                    f_dst.write(f_src.read())
            print(f"Copied file to slides folder: {slides_file_path}")

            # Process based on file type
            if file.filename.lower().endswith('.pdf'):
                # Process PDF file with both standard and enhanced extraction
                slides = extract_text_from_pdf(original_file_path)
                if not slides:
                    continue  # Skip if extraction fails
                    
                # Save basic extracted text to JSON for backward compatibility
                filename = file.filename.rsplit('.', 1)[0]
                save_extracted_pdf_text(slides, filename)
                print(f"Saved basic PDF extraction for: {filename}")
                
                # Save enhanced PDF extraction with visual elements
                enhanced_data = save_enhanced_pdf_extraction(original_file_path, filename)
                print(f"Saved enhanced PDF extraction: {enhanced_data is not None}")
                
                # If enhanced extraction was successful, use its slides (which contain more details)
                if enhanced_data and enhanced_data.get("slides"):
                    slides = enhanced_data.get("slides")
                
            elif file.filename.lower().endswith('.ppt'):
                # Convert .ppt to .pptx if needed
                converted_path = convert_ppt_to_pptx(original_file_path)
                if not converted_path:
                    continue  # Skip if conversion fails
                
                # Extract text from PowerPoint
                slides = extract_text_from_pptx(converted_path)
                
                # Save extracted text to JSON
                save_extracted_text(slides, file.filename.rsplit('.', 1)[0])
                
            else:  # .pptx file
                # Extract text from PowerPoint
                slides = extract_text_from_pptx(original_file_path)
                
                # Save extracted text to JSON
                save_extracted_text(slides, file.filename.rsplit('.', 1)[0])

            # Store this presentation's data
            presentation_data = {
                "filename": file.filename,  # Keep the full filename with extension
                "basename": file.filename.rsplit('.', 1)[0],  # Store the base name without extension
                "slides": slides,
                "file_type": file.filename.split('.')[-1].lower(),
                "original_path": original_file_path,
                "has_enhanced_data": file.filename.lower().endswith('.pdf')  # Only PDFs have enhanced data currently
            }
            
            presentations.append(presentation_data)

        if not presentations:
            return jsonify({"error": "No valid presentations uploaded"}), 400

        return jsonify({"presentations": presentations})
        
    except Exception as e:
        print(f"Error in upload_slide: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": f"Server error: {str(e)}"}), 500
    
@app.route('/generate-image', methods=['POST'])
def generate_image():
    """Generate images using Stability AI API."""
    # Get Stability AI API key from environment
    STABILITY_API_KEY = os.getenv('STABILITY_API_KEY')
    
    if not STABILITY_API_KEY:
        return jsonify({"error": "Stability AI API key not configured. Set STABILITY_API_KEY in .env file."}), 500
        
    try:
        data = request.get_json()
        prompt = data.get("prompt")

        if not prompt:
            return jsonify({"error": "Prompt is required"}), 400

        # Add detailed logging
        print(f"Sending request to Stability AI API with prompt: {prompt[:50]}...")
        
        # Call the Stability AI API with the correct endpoint and format
        try:
            response = requests.post(
                "https://api.stability.ai/v1/generation/stable-diffusion-xl-1024-v1-0/text-to-image",
                headers={
                    "Authorization": f"Bearer {STABILITY_API_KEY}",
                    "Content-Type": "application/json",
                    "Accept": "application/json"
                },
                json={
                    "text_prompts": [
                        {
                            "text": prompt,
                            "weight": 1
                        }
                    ],
                    "cfg_scale": 7,
                    "height": 1024,
                    "width": 1024,
                    "samples": 1,
                    "steps": 30
                },
                timeout=60  # Add a timeout
            )
            
            # Log the response status
            print(f"Stability AI API Response status: {response.status_code}")
            
            if response.status_code != 200:
                error_detail = response.json() if response.headers.get('content-type') == 'application/json' else response.text
                print(f"Error response: {error_detail}")
                return jsonify({
                    "error": "Failed to generate image", 
                    "details": error_detail,
                    "status_code": response.status_code
                }), 500
            
            # Parse the JSON response and get the base64 image data
            response_data = response.json()
            if "artifacts" in response_data and len(response_data["artifacts"]) > 0:
                img_data = response_data["artifacts"][0]["base64"]
                return jsonify({
                    "prompt": prompt,
                    "image_base64": f"data:image/png;base64,{img_data}"
                })
            else:
                return jsonify({"error": "No image generated in response"}), 500
            
        except requests.RequestException as e:
            print(f"Request exception: {str(e)}")
            return jsonify({"error": f"Request failed: {str(e)}"}), 500

    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"error": f"Server error: {str(e)}"}), 500

@app.route('/static/<path:path>')
def send_static(path):
    return send_from_directory('static', path)

@app.route('/')
def index():
    return send_from_directory('templates', 'index.html')

if __name__ == '__main__':
    app.run(debug=True)
